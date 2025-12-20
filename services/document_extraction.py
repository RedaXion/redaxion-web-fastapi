"""
Document Extraction Service - Extract text from PDF, DOCX, PPTX files

Extracts text content from uploaded documents for use as context
in AI-generated exams.
"""

import io
from typing import Optional

# PDF extraction
try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    print("⚠️ PyPDF2 not available - PDF extraction disabled")

# DOCX extraction
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    print("⚠️ python-docx not available - DOCX extraction disabled")

# PPTX extraction
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx not available - PPTX extraction disabled")


# Maximum characters to extract (to fit in GPT-4 context window)
MAX_CONTEXT_CHARS = 50000


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file."""
    if not PDF_AVAILABLE:
        return "[Error: PDF extraction not available]"
    
    try:
        pdf_file = io.BytesIO(file_bytes)
        reader = PdfReader(pdf_file)
        
        text_parts = []
        for page_num, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(f"--- Página {page_num + 1} ---\n{page_text}")
        
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"❌ Error extracting PDF: {e}")
        return f"[Error extracting PDF: {str(e)}]"


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from a DOCX file."""
    if not DOCX_AVAILABLE:
        return "[Error: DOCX extraction not available]"
    
    try:
        docx_file = io.BytesIO(file_bytes)
        doc = Document(docx_file)
        
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return "\n".join(text_parts)
    except Exception as e:
        print(f"❌ Error extracting DOCX: {e}")
        return f"[Error extracting DOCX: {str(e)}]"


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from a PowerPoint file."""
    if not PPTX_AVAILABLE:
        return "[Error: PPTX extraction not available]"
    
    try:
        pptx_file = io.BytesIO(file_bytes)
        prs = Presentation(pptx_file)
        
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            if slide_texts:
                text_parts.append(f"--- Diapositiva {slide_num} ---\n" + "\n".join(slide_texts))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"❌ Error extracting PPTX: {e}")
        return f"[Error extracting PPTX: {str(e)}]"


def extract_text_from_file(filename: str, file_bytes: bytes) -> str:
    """
    Extract text from a file based on its extension.
    
    Args:
        filename: Original filename with extension
        file_bytes: Raw bytes of the file
    
    Returns:
        Extracted text content
    """
    filename_lower = filename.lower()
    
    if filename_lower.endswith('.pdf'):
        return extract_text_from_pdf(file_bytes)
    elif filename_lower.endswith('.docx') or filename_lower.endswith('.doc'):
        return extract_text_from_docx(file_bytes)
    elif filename_lower.endswith('.pptx') or filename_lower.endswith('.ppt'):
        return extract_text_from_pptx(file_bytes)
    elif filename_lower.endswith('.txt'):
        # Plain text - just decode
        try:
            return file_bytes.decode('utf-8')
        except:
            try:
                return file_bytes.decode('latin-1')
            except:
                return "[Error: Could not decode text file]"
    else:
        return f"[Unsupported file type: {filename}]"


def extract_context_from_files(files: list) -> str:
    """
    Extract and combine text from multiple files.
    
    Args:
        files: List of tuples (filename, file_bytes)
    
    Returns:
        Combined text, truncated to MAX_CONTEXT_CHARS
    """
    if not files:
        return ""
    
    all_text_parts = []
    
    for filename, file_bytes in files:
        print(f"📄 Extracting text from: {filename}")
        text = extract_text_from_file(filename, file_bytes)
        if text and not text.startswith("[Error"):
            all_text_parts.append(f"=== {filename} ===\n{text}")
    
    combined_text = "\n\n".join(all_text_parts)
    
    # Truncate if too long
    if len(combined_text) > MAX_CONTEXT_CHARS:
        combined_text = combined_text[:MAX_CONTEXT_CHARS] + "\n\n[... contenido truncado por límite de contexto ...]"
        print(f"⚠️ Context truncated to {MAX_CONTEXT_CHARS} characters")
    
    print(f"✅ Extracted {len(combined_text)} characters from {len(files)} file(s)")
    return combined_text
